import streamlit as st
import os

# Google Gemini API imports
import google.generativeai as genai

# File parsers
from PyPDF2 import PdfReader
from docx import Document
import openpyxl
from pptx import Presentation

# Text processing
from textblob import TextBlob

# Import the unified detector
from enhanced_unified_agent_detector import detect_agent

# Set your Google Gemini API key here
GEMINI_API_KEY = "AIzaSyCkoZ3ubDy7kSNwiD0dBgNNHftdDez0hYw"

# Set up Gemini API
genai.configure(api_key=GEMINI_API_KEY)

# Model config
GENAI_MODEL = 'gemini-2.0-flash-lite'

def gemini_prompt(prompt, temperature=0.2):
    model = genai.GenerativeModel(GENAI_MODEL)
    response = model.generate_content(prompt, generation_config={"temperature": temperature})
    return response.text

def emotion_analysis_agent(query):
    prompt = (
        "You are a close, emotionally intelligent friend who truly listens and cares. "
        "I am about to share something personal. Please focus fully on what I say, without distractions or unrelated content. "
        
        "Perform a comprehensive emotional analysis by examining ALL keywords and phrases that indicate emotional states. "
        "Consider not just direct emotion words, but also:\n"
        "- Synonyms and related expressions (e.g., 'devastated' = sad, 'ecstatic' = happy, 'furious' = angry)\n"
        "- Contextual emotional indicators (e.g., 'heartbroken', 'over the moon', 'at my wit's end')\n"
        "- Intensity modifiers (e.g., 'slightly', 'extremely', 'somewhat', 'utterly')\n"
        "- Body language descriptions (e.g., 'tears streaming', 'heart racing', 'shoulders slumped')\n"
        "- Metaphorical expressions (e.g., 'feeling empty', 'on cloud nine', 'carrying the weight')\n"
        "- Tone and sentiment indicators in phrasing\n"
        "- Subtle emotional cues in word choice and sentence structure\n\n"
        
        "Create a detailed emotional profile that includes:\n"
        "1. Primary emotions detected (with confidence levels)\n"
        "2. Secondary or underlying emotions\n"
        "3. Emotional intensity (scale of 1-10)\n"
        "4. Specific keywords/phrases that triggered each emotion\n"
        "5. Emotional complexity or mixed feelings\n"
        "6. Possible emotional needs or support required\n\n"
        
        "Reflect back with deep understanding and genuine empathy — not like a therapist or AI, but like a trusted friend who truly sees and hears you. "
        "Acknowledge the complexity of human emotions and validate whatever you're feeling without judgment. "
        "Here is what I want to share:\n\n{query}"
    )
    return gemini_prompt(prompt, temperature=0.1)

def maths_agent(query):
    prompt = (
        "You are a helpful math assistant. Solve the following math problem and show all steps. "
        "Recognize not just direct math terms (like 'solve', 'calculate', 'integral'), but also synonyms and related phrases that indicate a math question. "
        f"\n\n{query}"
    )
    return gemini_prompt(prompt)

def physics_chem_agent(query):
    prompt = (
        "You are a science assistant. Provide a detailed solution for this physics or chemistry problem. "
        "Recognize not just direct science terms (like 'physics', 'chemistry', 'reaction'), but also synonyms and related phrases that indicate a science question. "
        f"\n\n{query}"
    )
    return gemini_prompt(prompt)

def coding_agent(query):
    prompt = (
        "You are a coding assistant. Write and explain code for the following request. Provide working code and a short explanation. "
        "Recognize not just direct coding terms (like 'code', 'python', 'algorithm'), but also synonyms and related phrases that indicate a coding or programming request. "
        f"\n\n{query}"
    )
    return gemini_prompt(prompt)

def general_conversation_agent(query):
    prompt = (
        "You are a helpful, friendly assistant. Respond to the following input. "
        "Recognize not just direct conversation starters (like 'hello', 'hi'), but also synonyms and related phrases that indicate a general conversation or request for help. "
        f"\n\n{query}"
    )
    return gemini_prompt(prompt)

def extract_text_from_file(uploaded_file):
    ext = uploaded_file.name.split('.')[-1].lower()
    if ext == 'pdf':
        reader = PdfReader(uploaded_file)
        text = "".join(page.extract_text() for page in reader.pages)
        return text
    elif ext == 'docx':
        doc = Document(uploaded_file)
        return "\n".join([para.text for para in doc.paragraphs])
    elif ext in ('xlsx', 'xls'):
        wb = openpyxl.load_workbook(uploaded_file)
        text = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                text.append('\t'.join([str(cell) for cell in row if cell]))
        return '\n'.join(text)
    elif ext == 'pptx':
        prs = Presentation(uploaded_file)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    else:
        return "Unsupported file format"

# Map agent names to their functions
AGENTS = {
    'Emotion Analysis': emotion_analysis_agent,
    'Math': maths_agent,
    'Physics/Chemistry': physics_chem_agent,
    'Coding': coding_agent,
    'General Conversation': general_conversation_agent
}

st.title("Multi-Agent Gemini Assistant")
st.write("Powered by Gemini-2.0-Flash-Lite via Google AI Studio")

uploaded_file = st.file_uploader("Upload file (.pdf, .docx, .xlsx, .pptx)", type=['pdf', 'docx', 'xlsx', 'xls', 'pptx'])
user_query = st.text_area("Or enter your question/task here:")

if uploaded_file:
    st.info("Extracting text from uploaded file...")
    file_text = extract_text_from_file(uploaded_file)
    st.text_area("Extracted File Text", file_text, height=200)
else:
    file_text = ""

if st.button("Clear Chat History"):
    st.session_state['history'] = []
if 'history' not in st.session_state:
    st.session_state['history'] = []

if st.button("Run"):
    input_text = user_query.strip()
    if not input_text and not file_text:
        st.warning("Please provide a question or upload a file.")
    else:
        combined_text = f"Context:\n{file_text}\n\nQuestion:\n{input_text}" if file_text else input_text
        agent_choice = detect_agent(combined_text)
        with st.spinner(f"Using agent: {agent_choice} ..."):
            try:
                result = AGENTS[agent_choice](combined_text)
                st.session_state['history'].append({
                    'agent': agent_choice,
                    'input': combined_text,
                    'output': result
                })
            except Exception as e:
                st.error(f"Error running agent: {e}")

# Display chat history
for i, chat in enumerate(st.session_state['history']):
    st.markdown(f"**[{chat['agent']}] Input:** {chat['input']}")
    st.markdown(f"**[{chat['agent']}] Output:**\n{chat['output']}")
    st.markdown("---")

# Sample initialization for testing
if 'history' not in st.session_state:
    st.session_state['history'] = [
        {'agent': 'Coding', 'input': 'Write Python function', 'output': 'def foo(): pass'},
        {'agent': 'Math', 'input': 'Solve x^2=4', 'output': 'x=2 or x=-2'}
    ]




with st.sidebar:
    st.header("Chat History")
    with st.expander("Show/Hide Chat History", expanded=False):
        if st.session_state.get('history'):
            for chat in st.session_state['history']:
                input_html = chat['input'].replace('\n', '<br>')
                output_html = chat['output'].replace('\n', '<br>')
                with st.container():
                    st.markdown(
                        f"""
                        <div style="
                            border: 1px solid #ddd;
                            padding: 10px;
                            margin-bottom: 10px;
                            border-radius: 5px;
                            background-color: black;
                        ">
                        <strong>[{chat['agent']}] Input:</strong><br>{input_html}<br><br>
                        <strong>[{chat['agent']}] Output:</strong><br>{output_html}
                        </div>
                        """, 
                        unsafe_allow_html=True
                    )
        else:
            st.info("No chat history available.")
